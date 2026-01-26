"""Export Service - Converts HTML/SlideSpec to PPTX and DOCX."""

import io
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

from app.schemas.slidespec import SlideSpec, Slide, Element


class PPTXExporter:
    """Exports SlideSpec/HTML to PPTX format."""

    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Margins
    MARGIN_LEFT = Inches(0.5)
    MARGIN_RIGHT = Inches(0.5)
    MARGIN_TOP = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)

    # Colors
    PRIMARY_COLOR = RGBColor(0x1a, 0x1a, 0x2e)
    ACCENT_COLOR = RGBColor(0x4a, 0x90, 0xd9)
    TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
    LIGHT_TEXT = RGBColor(0x66, 0x66, 0x66)

    def __init__(self, template_path: Path | str | None = None):
        """Initialize with optional template."""
        self.template_path = template_path

    def create_presentation(self) -> Presentation:
        """Create a new presentation."""
        if self.template_path and Path(self.template_path).exists():
            prs = Presentation(str(self.template_path))
        else:
            prs = Presentation()
            # Set slide dimensions to 16:9
            prs.slide_width = self.SLIDE_WIDTH
            prs.slide_height = self.SLIDE_HEIGHT
        return prs

    def export_slidespec(self, slidespec: SlideSpec) -> bytes:
        """Export SlideSpec to PPTX bytes."""
        prs = self.create_presentation()

        for slide_data in slidespec.slides:
            self._add_slide(prs, slide_data)

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def _add_slide(self, prs: Presentation, slide: Slide):
        """Add a slide to the presentation."""
        # Use blank layout for custom positioning
        blank_layout = prs.slide_layouts[6]  # Blank
        ppt_slide = prs.slides.add_slide(blank_layout)

        layout_id = slide.layout.layout_id if slide.layout else "one_column"

        if layout_id == "title_center":
            self._render_title_slide(ppt_slide, slide)
        elif layout_id == "section_header":
            self._render_section_slide(ppt_slide, slide)
        elif layout_id == "two_column":
            self._render_two_column_slide(ppt_slide, slide)
        elif layout_id in ("chart_focus", "table_focus"):
            self._render_content_slide(ppt_slide, slide)
        elif layout_id == "quote_center":
            self._render_quote_slide(ppt_slide, slide)
        elif layout_id == "closing":
            self._render_closing_slide(ppt_slide, slide)
        else:
            self._render_content_slide(ppt_slide, slide)

        # Add speaker notes if present
        if slide.speaker_notes:
            notes_slide = ppt_slide.notes_slide
            notes_slide.notes_text_frame.text = slide.speaker_notes

    def _render_title_slide(self, ppt_slide, slide: Slide):
        """Render a title slide."""
        title_text = ""
        subtitle_text = ""

        for elem in slide.elements:
            if elem.role == "title" or (elem.kind == "text" and not title_text):
                title_text = elem.content.get("text", "")
            elif elem.role == "subtitle" or elem.kind == "text":
                subtitle_text = elem.content.get("text", "")

        # Title
        if title_text:
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(11.333)
            height = Inches(1.5)
            txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle_text:
            left = Inches(1.5)
            top = Inches(4.2)
            width = Inches(10.333)
            height = Inches(1)
            txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.font.size = Pt(24)
            p.font.color.rgb = self.LIGHT_TEXT
            p.alignment = PP_ALIGN.CENTER

    def _render_section_slide(self, ppt_slide, slide: Slide):
        """Render a section header slide."""
        title_text = ""
        subtitle_text = ""

        for elem in slide.elements:
            if elem.role == "title" or (elem.kind == "text" and not title_text):
                title_text = elem.content.get("text", "")
            elif elem.kind == "text":
                subtitle_text = elem.content.get("text", "")

        # Title with accent bar
        if title_text:
            # Accent bar
            left = Inches(1)
            top = Inches(2.8)
            width = Inches(0.1)
            height = Inches(1.5)
            shape = ppt_slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.ACCENT_COLOR
            shape.line.fill.background()

            # Title text
            left = Inches(1.3)
            top = Inches(2.8)
            width = Inches(10)
            height = Inches(1.2)
            txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(42)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR

        if subtitle_text:
            left = Inches(1.3)
            top = Inches(4.2)
            width = Inches(10)
            height = Inches(0.8)
            txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.font.size = Pt(20)
            p.font.color.rgb = self.LIGHT_TEXT

    def _render_content_slide(self, ppt_slide, slide: Slide):
        """Render a standard content slide."""
        title_elem = None
        content_elems = []

        for elem in slide.elements:
            if elem.role == "title":
                title_elem = elem
            else:
                content_elems.append(elem)

        # Title
        if title_elem:
            left = Inches(0.5)
            top = Inches(0.4)
            width = Inches(12.333)
            height = Inches(0.8)
            txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title_elem.content.get("text", "")
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY_COLOR

        # Content
        top_offset = Inches(1.4)
        for elem in content_elems:
            if elem.kind == "text":
                left = Inches(0.5)
                width = Inches(12.333)
                height = Inches(1)
                txBox = ppt_slide.shapes.add_textbox(left, top_offset, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = elem.content.get("text", "")
                p.font.size = Pt(18)
                p.font.color.rgb = self.TEXT_COLOR
                top_offset += Inches(1.2)

            elif elem.kind == "bullets":
                items = elem.content.get("items", [])
                left = Inches(0.5)
                width = Inches(12.333)
                height = Inches(len(items) * 0.5 + 0.5)
                txBox = ppt_slide.shapes.add_textbox(left, top_offset, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                for i, item in enumerate(items):
                    text = item if isinstance(item, str) else item.get("text", "")
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {text}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = self.TEXT_COLOR
                    p.space_after = Pt(8)

                top_offset += height

            elif elem.kind == "table":
                self._add_table(ppt_slide, elem, top_offset)
                rows = len(elem.content.get("rows", []))
                top_offset += Inches(0.5 + rows * 0.4)

    def _render_two_column_slide(self, ppt_slide, slide: Slide):
        """Render a two-column slide."""
        # Simplified - just use content slide layout
        self._render_content_slide(ppt_slide, slide)

    def _render_quote_slide(self, ppt_slide, slide: Slide):
        """Render a quote slide."""
        quote_text = ""
        source_text = ""

        for elem in slide.elements:
            if elem.role == "quote" or (elem.kind == "text" and not quote_text):
                quote_text = elem.content.get("text", "")
            elif elem.role in ("source", "attribution"):
                source_text = elem.content.get("text", "")

        if quote_text:
            left = Inches(1.5)
            top = Inches(2.5)
            width = Inches(10.333)
            height = Inches(2)
            txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f'"{quote_text}"'
            p.font.size = Pt(28)
            p.font.italic = True
            p.font.color.rgb = self.PRIMARY_COLOR
            p.alignment = PP_ALIGN.CENTER

        if source_text:
            left = Inches(1.5)
            top = Inches(4.8)
            width = Inches(10.333)
            height = Inches(0.6)
            txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {source_text}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.LIGHT_TEXT
            p.alignment = PP_ALIGN.CENTER

    def _render_closing_slide(self, ppt_slide, slide: Slide):
        """Render a closing slide."""
        title_text = "Thank You"

        for elem in slide.elements:
            if elem.role == "title" or elem.kind == "text":
                title_text = elem.content.get("text", title_text)
                break

        left = Inches(1)
        top = Inches(3)
        width = Inches(11.333)
        height = Inches(1.5)
        txBox = ppt_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER

    def _add_table(self, ppt_slide, elem: Element, top: float):
        """Add a table to a slide."""
        content = elem.content
        columns = content.get("columns", [])
        rows = content.get("rows", [])

        if not columns or not rows:
            return

        left = Inches(0.5)
        width = Inches(12.333)
        row_height = Inches(0.4)
        height = row_height * (len(rows) + 1)

        table = ppt_slide.shapes.add_table(
            len(rows) + 1, len(columns), left, top, width, height
        ).table

        # Header row
        for i, col in enumerate(columns):
            cell = table.cell(0, i)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.PRIMARY_COLOR
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(12)
            p.font.bold = True

        # Data rows
        for i, row in enumerate(rows):
            for j, cell_value in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(cell_value) if cell_value is not None else ""
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = self.TEXT_COLOR


class ImageBasedPPTXExporter:
    """Exports HTML slides to PPTX as images for pixel-perfect rendering."""

    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self):
        pass

    def create_presentation(self) -> Presentation:
        """Create a new presentation with 16:9 aspect ratio."""
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT
        return prs

    def export_from_images(self, images: list[bytes], speaker_notes: list[str] | None = None) -> bytes:
        """Export a list of PNG image bytes to PPTX.

        Args:
            images: List of PNG image bytes, one per slide
            speaker_notes: Optional list of speaker notes for each slide
        """
        prs = self.create_presentation()
        blank_layout = prs.slide_layouts[6]  # Blank layout

        for idx, image_bytes in enumerate(images):
            slide = prs.slides.add_slide(blank_layout)

            # Add image as background (full slide)
            image_stream = io.BytesIO(image_bytes)
            left = Inches(0)
            top = Inches(0)
            width = self.SLIDE_WIDTH
            height = self.SLIDE_HEIGHT

            slide.shapes.add_picture(image_stream, left, top, width, height)

            # Add speaker notes if provided
            if speaker_notes and idx < len(speaker_notes) and speaker_notes[idx]:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = speaker_notes[idx]

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()


class DOCXExporter:
    """Exports SlideSpec/HTML to DOCX format."""

    def __init__(self, template_path: Path | str | None = None):
        """Initialize with optional template."""
        self.template_path = template_path

    def export_slidespec(self, slidespec: SlideSpec) -> bytes:
        """Export SlideSpec to DOCX bytes (as a document outline)."""
        doc = Document()

        # Title
        doc.add_heading(slidespec.deck.title, 0)

        if slidespec.deck.subtitle:
            doc.add_paragraph(slidespec.deck.subtitle)

        # Add each slide as a section
        for slide in slidespec.slides:
            self._add_slide_section(doc, slide)

        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def _add_slide_section(self, doc: Document, slide: Slide):
        """Add a slide as a document section."""
        # Find title
        title_text = slide.title or ""
        for elem in slide.elements:
            if elem.role == "title":
                title_text = elem.content.get("text", "")
                break

        if title_text:
            level = 1 if slide.type == "section" else 2
            doc.add_heading(title_text, level)

        # Add content
        for elem in slide.elements:
            if elem.role == "title":
                continue

            if elem.kind == "text":
                doc.add_paragraph(elem.content.get("text", ""))

            elif elem.kind == "bullets":
                items = elem.content.get("items", [])
                for item in items:
                    text = item if isinstance(item, str) else item.get("text", "")
                    doc.add_paragraph(text, style="List Bullet")

            elif elem.kind == "table":
                self._add_table(doc, elem)

        doc.add_paragraph()  # Spacing

    def _add_table(self, doc: Document, elem: Element):
        """Add a table to the document."""
        content = elem.content
        columns = content.get("columns", [])
        rows = content.get("rows", [])

        if not columns or not rows:
            return

        table = doc.add_table(rows=len(rows) + 1, cols=len(columns))
        table.style = "Table Grid"

        # Header
        for i, col in enumerate(columns):
            cell = table.rows[0].cells[i]
            cell.text = str(col)
            for p in cell.paragraphs:
                for run in p.runs:
                    run.bold = True

        # Data
        for i, row in enumerate(rows):
            for j, value in enumerate(row):
                table.rows[i + 1].cells[j].text = str(value) if value else ""


class ExportService:
    """Unified export service for PPTX and DOCX."""

    def __init__(self):
        self.pptx_exporter = PPTXExporter()
        self.image_pptx_exporter = ImageBasedPPTXExporter()
        self.docx_exporter = DOCXExporter()

    def export_to_pptx(self, slidespec: SlideSpec) -> bytes:
        """Export SlideSpec to PPTX (element-based, legacy)."""
        return self.pptx_exporter.export_slidespec(slidespec)

    async def export_to_pptx_as_images(self, slidespec: SlideSpec) -> bytes:
        """Export SlideSpec to PPTX as images (pixel-perfect rendering).

        This captures each slide as an image and embeds them into PPTX,
        ensuring the output looks exactly like the HTML preview.
        """
        from app.services.html_capture_service import get_html_capture_service

        capture_service = get_html_capture_service()

        try:
            # Capture all slides as images
            images = await capture_service.capture_slidespec(slidespec)

            # Collect speaker notes
            speaker_notes = [
                slide.speaker_notes for slide in slidespec.slides
            ]

            # Create PPTX from images
            return self.image_pptx_exporter.export_from_images(images, speaker_notes)
        finally:
            # Don't close the service here - it's a singleton
            pass

    async def export_html_slides_to_pptx(self, html_slides: list[str], speaker_notes: list[str] | None = None) -> bytes:
        """Export HTML slides to PPTX as images.

        Args:
            html_slides: List of HTML strings, one per slide
            speaker_notes: Optional list of speaker notes
        """
        from app.services.html_capture_service import get_html_capture_service

        capture_service = get_html_capture_service()

        try:
            # Capture all HTML slides as images
            images = await capture_service.capture_html_slides(html_slides)

            # Create PPTX from images
            return self.image_pptx_exporter.export_from_images(images, speaker_notes)
        finally:
            pass

    def export_to_docx(self, slidespec: SlideSpec) -> bytes:
        """Export SlideSpec to DOCX."""
        return self.docx_exporter.export_slidespec(slidespec)

    def html_to_pptx(self, html_slides: list[str], slidespec: SlideSpec | None = None) -> bytes:
        """Convert HTML slides to PPTX (uses SlideSpec if available, legacy)."""
        if slidespec:
            return self.export_to_pptx(slidespec)
        raise NotImplementedError("Use export_html_slides_to_pptx for HTML-only export")

    def html_to_docx(self, html_content: str, slidespec: SlideSpec | None = None) -> bytes:
        """Convert HTML to DOCX (uses SlideSpec if available)."""
        if slidespec:
            return self.export_to_docx(slidespec)
        raise NotImplementedError("HTML-only export not yet implemented")
